"""
DualTrack API: 声明式 + 精确控制的双轨设计

核心原则：
1. 轨道 1（声明式）：面向 LLM，用意图描述代替坐标计算
2. 轨道 2（精确控制）：面向人类微调，可覆盖任何自动计算的值
3. 轨道 3（底层原生）：完全绕过声明式，直接写 python-pptx 代码

使用示例见文件底部 __main__ 部分。
"""

from dataclasses import dataclass, field
from typing import List, Dict, Optional, Tuple, Union, Any
from enum import Enum


@dataclass
class Position:
    """精确位置（英寸）"""
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    
    def override(self, left=None, top=None, width=None, height=None):
        """创建覆盖后的新 Position"""
        return Position(
            left=left if left is not None else self.left,
            top=top if top is not None else self.top,
            width=width if width is not None else self.width,
            height=height if height is not None else self.height,
        )


@dataclass 
class TypographySpec:
    """字体规格"""
    font_name: str = "Microsoft YaHei"
    font_size: float = 9.5
    color: str = "#1A1A1A"
    bold: bool = False
    italic: bool = False
    line_spacing: float = 1.45
    
    def override(self, **kwargs):
        """创建覆盖后的新 TypographySpec"""
        data = {
            "font_name": kwargs.get("font_name", self.font_name),
            "font_size": kwargs.get("font_size", self.font_size),
            "color": kwargs.get("color", self.color),
            "bold": kwargs.get("bold", self.bold),
            "italic": kwargs.get("italic", self.italic),
            "line_spacing": kwargs.get("line_spacing", self.line_spacing),
        }
        return TypographySpec(**data)


class DeclarativeElement:
    """
    声明式元素基类
    
    所有声明式元素都支持：
    - 自动计算默认属性（由布局引擎在编译时填充）
    - 显式覆盖任何属性（保留精细控制）
    - 导出为底层 python-pptx 代码
    """
    
    def __init__(self, element_type: str):
        self.element_type = element_type
        
        # 自动计算的值（由引擎填充）
        self._auto_position: Optional[Position] = None
        self._auto_typography: Optional[TypographySpec] = None
        self._auto_color: Optional[str] = None
        
        # 显式覆盖的值（用户/LLM 设置）
        self._override_position: Optional[Position] = None
        self._override_typography: Optional[TypographySpec] = None
        self._override_color: Optional[str] = None
        
        # 编译后的最终值（auto + override 合并）
        self._compiled = False
        self._final_position: Optional[Position] = None
        self._final_typography: Optional[TypographySpec] = None
        self._final_color: Optional[str] = None
    
    # ═══ 轨道 1：声明式接口（描述意图）═══
    
    def set_semantic_color(self, role: str):
        """
        声明式设置语义颜色
        
        Args:
            role: "positive" | "warning" | "contrast" | "neutral"
        """
        # 由引擎在编译时解析为具体 hex
        self._semantic_color_role = role
        return self
    
    def set_content(self, text: str):
        """声明式设置内容文本"""
        self.text = text
        return self
    
    # ═══ 轨道 2：精确覆盖接口（精细调整）═══
    
    def override_position(self, left=None, top=None, width=None, height=None):
        """
        精确覆盖位置（英寸）
        
        示例：
            card.override_position(height=1.8)  # 只改高度，其他保持自动计算
            card.override_position(left=0.5, top=1.2, width=4.0, height=2.0)  # 全改
        """
        self._override_position = Position(left or 0, top or 0, width or 0, height or 0)
        self._compiled = False
        return self
    
    def override_typography(self, font_name=None, font_size=None, color=None, 
                           bold=None, italic=None, line_spacing=None):
        """精确覆盖字体样式"""
        spec = TypographySpec()
        if font_name: spec.font_name = font_name
        if font_size: spec.font_size = font_size
        if color: spec.color = color
        if bold is not None: spec.bold = bold
        if italic is not None: spec.italic = italic
        if line_spacing: spec.line_spacing = line_spacing
        self._override_typography = spec
        self._compiled = False
        return self
    
    def override_color(self, hex_color: str):
        """精确覆盖颜色（hex 值）"""
        self._override_color = hex_color
        self._compiled = False
        return self
    
    # ═══ 编译：合并 auto + override ═══
    
    def compile(self, design_contract=None):
        """
        编译最终值
        
        合并规则：
        - 如果 override 存在，使用 override
        - 否则使用 auto（由引擎计算）
        - 如果两者都不存在，使用默认值
        """
        # 位置合并
        base_pos = self._auto_position or Position()
        if self._override_position:
            self._final_position = base_pos.override(
                left=self._override_position.left or None,
                top=self._override_position.top or None,
                width=self._override_position.width or None,
                height=self._override_position.height or None,
            )
        else:
            self._final_position = base_pos
        
        # 字体合并
        base_typo = self._auto_typography or TypographySpec()
        if self._override_typography:
            self._final_typography = base_typo.override(
                font_name=self._override_typography.font_name,
                font_size=self._override_typography.font_size,
                color=self._override_typography.color,
                bold=self._override_typography.bold,
                italic=self._override_typography.italic,
                line_spacing=self._override_typography.line_spacing,
            )
        else:
            self._final_typography = base_typo
        
        # 颜色合并
        if self._override_color:
            self._final_color = self._override_color
        elif self._auto_color:
            self._final_color = self._auto_color
        elif hasattr(self, '_semantic_color_role') and design_contract:
            self._final_color = design_contract.get_color(self._semantic_color_role)
        else:
            self._final_color = "#1A1A1A"
        
        self._compiled = True
        return self
    
    # ═══ 导出：生成底层代码 ═══
    
    def to_python_code(self) -> str:
        """
        导出为精确的 python-pptx 代码
        
        这是从声明式到底层的最终编译步骤。
        生成的代码可以直接执行，产生像素级精确的结果。
        """
        if not self._compiled:
            raise RuntimeError("必须先调用 compile() 再导出代码")
        
        p = self._final_position
        t = self._final_typography
        
        # 生成底层代码（示例）
        code = f"""
# 自动生成的精确代码（from {self.element_type}）
{self._generate_shape_code()}
"""
        return code
    
    def _generate_shape_code(self) -> str:
        """由子类实现"""
        return "# 基类不生成代码"
    
    # ═══ 属性访问器 ═══
    
    @property
    def position(self) -> Position:
        if not self._compiled:
            self.compile()
        return self._final_position
    
    @property
    def typography(self) -> TypographySpec:
        if not self._compiled:
            self.compile()
        return self._final_typography
    
    @property
    def color(self) -> str:
        if not self._compiled:
            self.compile()
        return self._final_color


class CardElement(DeclarativeElement):
    """内容卡片元素"""
    
    def __init__(self, title: str = "", body: str = ""):
        super().__init__("card")
        self.title = title
        self.body = body
        self.accent_color_role = "neutral"
    
    def set_accent(self, role: str):
        """声明式设置强调线语义颜色"""
        self.accent_color_role = role
        return self
    
    def _generate_shape_code(self) -> str:
        p = self._final_position
        t = self._final_typography
        
        return f"""
# Card: {self.title[:20]}...
_rect(sl, {p.left}, {p.top}, {p.width}, {p.height}, WHT)
_line(sl, {p.left}, {p.top}, {p.width}, {self._final_color}, 0.03)
if "{self.title}":
    tx(sl, {p.left + 0.15}, {p.top + 0.08}, {p.width - 0.28}, 0.26, 
       "{self.title}", fs=12, bd=True, cl=NAV)
if "{self.body}":
    tx(sl, {p.left + 0.15}, {p.top + 0.36}, {p.width - 0.28}, {p.height - 0.42},
       "{self.body}", fs={t.font_size}, cl=GRY, ls={t.line_spacing})
"""


class HeroElement(DeclarativeElement):
    """Hero 统计元素"""
    
    def __init__(self, number: str = "", label: str = ""):
        super().__init__("hero")
        self.number = number
        self.label = label
    
    def _generate_shape_code(self) -> str:
        p = self._final_position
        return f"""
# Hero: {self.number} {self.label[:15]}...
hero(sl, {p.left}, {p.top}, {p.width}, {p.height}, 
     "{self.number}", "{self.label}", accent={self._final_color})
"""


class ImageElement(DeclarativeElement):
    """图片元素"""
    
    def __init__(self, image_path: str = ""):
        super().__init__("image")
        self.image_path = image_path
    
    def _generate_shape_code(self) -> str:
        p = self._final_position
        return f"""
# Image: {self.image_path}
img_fit(sl, "{self.image_path}", {p.left}, {p.top}, {p.width}, {p.height}, align="center")
"""


class RawElement:
    """
    轨道 3：底层原生元素
    
    完全绕过声明式系统，直接写 python-pptx 代码。
    用于完全自定义的场景。
    """
    
    def __init__(self, code_generator: callable):
        self.code_generator = code_generator
    
    def to_python_code(self) -> str:
        return self.code_generator()


# ═══════════════════════════════════════════
# Slide 类：声明式页面容器
# ═══════════════════════════════════════════

class DeclarativeSlide:
    """
    声明式幻灯片
    
    支持：
    - 声明式添加元素（自动计算位置）
    - 精确覆盖任何元素的属性
    - 插入底层原生代码
    """
    
    def __init__(self, layout_type: str):
        self.layout_type = layout_type
        self.elements: List[Union[DeclarativeElement, RawElement]] = []
        self.title: str = ""
        self.subtitle: str = ""
    
    def set_title(self, title: str, subtitle: str = ""):
        self.title = title
        self.subtitle = subtitle
        return self
    
    def add_card(self, title: str = "", body: str = "", accent: str = "neutral") -> CardElement:
        """声明式添加卡片"""
        card = CardElement(title, body)
        card.set_accent(accent)
        self.elements.append(card)
        return card  # 返回引用，支持链式覆盖
    
    def add_hero(self, number: str, label: str) -> HeroElement:
        """声明式添加 Hero 统计"""
        hero = HeroElement(number, label)
        self.elements.append(hero)
        return hero
    
    def add_image(self, path: str) -> ImageElement:
        """声明式添加图片"""
        img = ImageElement(path)
        self.elements.append(img)
        return img
    
    def add_raw(self, code_generator: callable):
        """添加底层原生代码"""
        self.elements.append(RawElement(code_generator))
        return self
    
    def compile(self, design_contract=None):
        """编译所有元素"""
        for elem in self.elements:
            if isinstance(elem, DeclarativeElement):
                elem.compile(design_contract)
        return self
    
    def to_python_code(self) -> str:
        """导出完整页面的 Python 代码"""
        lines = [f"# === Slide: {self.title} ==="]
        
        for elem in self.elements:
            lines.append(elem.to_python_code())
        
        return "\n".join(lines)


# ═══════════════════════════════════════════
# SlideDeck 类：顶层容器
# ═══════════════════════════════════════════

class SlideDeck:
    """
    声明式 PPT 容器
    
    双轨 API 入口：
    1. 快速声明：deck.add_slide(layout="...", title="...", ...)
    2. 精细调整：slide.add_card(...).override_position(height=1.8)
    3. 完全自定义：slide.add_raw(lambda: "底层代码")
    """
    
    def __init__(self, theme: str = "academic"):
        from design_contract import DesignContract, ThemeType
        self.contract = DesignContract(ThemeType(theme))
        self.slides: List[DeclarativeSlide] = []
    
    def add_slide(self, layout: str, title: str = "", subtitle: str = "") -> DeclarativeSlide:
        """添加新幻灯片"""
        slide = DeclarativeSlide(layout)
        slide.set_title(title, subtitle)
        self.slides.append(slide)
        return slide
    
    def build(self, output_path: str):
        """
        构建完整 deck
        
        流程：
        1. 编译所有 slide（自动计算坐标）
        2. 全局一致性检查
        3. 导出为 Python 代码
        4. 执行生成 PPTX
        """
        # 编译
        for slide in self.slides:
            slide.compile(self.contract)
        
        # 全局检查
        is_pass, violations = self.contract.validate_global()
        if not is_pass:
            print("全局一致性检查失败:")
            for v in violations:
                print(f"  [{v['severity']}] {v['message']}")
        
        # 导出代码（示例）
        code_parts = ["from pptx import Presentation", "prs = Presentation()", ""]
        for slide in self.slides:
            code_parts.append(slide.to_python_code())
            code_parts.append("")
        code_parts.append(f'prs.save("{output_path}")')
        
        full_code = "\n".join(code_parts)
        
        # 实际执行（这里仅演示）
        print(f"=== 生成的 Python 代码（{len(full_code)} 字符）===")
        print(full_code[:500] + "...")
        
        return full_code
    
    def get_contract_summary(self) -> str:
        """获取设计契约摘要（供 LLM 参考）"""
        return self.contract.get_context_summary()


# ═══════════════════════════════════════════
# 使用示例
# ═══════════════════════════════════════════

if __name__ == "__main__":
    
    # 示例 1：纯声明式（90% 场景）
    deck = SlideDeck(theme="medical")
    
    slide1 = deck.add_slide("cover", title="基于深度学习的医学影像分析")
    # 封面由引擎自动处理，无需添加元素
    
    slide2 = deck.add_slide("two_column_compare", title="研究背景与问题")
    slide2.add_card(
        title="现有方法",
        body="依赖人工特征提取\n泛化能力差\n分割精度不足",
        accent="warning"
    )
    slide2.add_card(
        title="本文方法",
        body="端到端深度学习\n引入注意力机制\nGrad-CAM 可解释性",
        accent="positive"
    )
    
    slide3 = deck.add_slide("timeline", title="方法流程")
    # 时间线元素由引擎根据内容自动布局
    
    slide4 = deck.add_slide("hero_panel", title="实验结果")
    slide4.add_hero("93.3%", "Dice 系数")
    slide4.add_hero("4.7%", "性能提升")
    
    # 示例 2：声明式 + 精确覆盖（9% 场景）
    slide5 = deck.add_slide("content_single", title="详细分析")
    card = slide5.add_card(
        title="消融实验",
        body="移除注意力模块后，Dice 系数下降至 89.1%..."
    )
    # 自动计算高度不够？直接覆盖
    card.override_position(height=2.0)  # 改为 2.0 英寸
    card.override_typography(font_size=10)  # 字号稍大
    
    # 示例 3：完全自定义（1% 场景）
    slide6 = deck.add_slide("custom", title="特殊图表")
    slide6.add_raw(lambda: """
# 完全自定义的复杂图表
from pptx.enum.chart import XL_CHART_TYPE
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
    Inches(8), Inches(5), chart_data
).chart
""")
    
    # 构建
    deck.build("output.pptx")
    
    # 查看设计契约摘要
    print("\n=== 设计契约摘要 ===")
    print(deck.get_contract_summary())
